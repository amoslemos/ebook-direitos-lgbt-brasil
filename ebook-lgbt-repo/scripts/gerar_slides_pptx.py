from pptx import Presentation
from pptx.util import Inches, Cm
import os

ASSETS = 'assets'
OUTPUT = 'ebook_lgbt_slides.pptx'

prs = Presentation()
prs.slide_width = Cm(33.87)
prs.slide_height = Cm(19.05)

imgs = [
    'capa.png',
    'contracapa.png',
    'infografico_espectro.png',
    'infografico_linha_do_tempo.png',
    'infografico_denuncia.png'
]

for img in imgs:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    path = os.path.join(ASSETS, img)
    if os.path.exists(path):
        pic = slide.shapes.add_picture(path, Cm(0), Cm(0), width=Cm(33.87), height=Cm(19.05))

prs.save(OUTPUT)
print('Slides gerados:', OUTPUT)
